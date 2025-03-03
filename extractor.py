import os
import logging
from abc import ABC, abstractmethod
from pptx import Presentation
from typing import Union, List, Dict

# Configure logging for better debugging and traceability.
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)

class BaseExtractor(ABC):
    """
    Abstract base class for file extractors.
    """

    @staticmethod
    @abstractmethod
    def extract(file_path: str) -> Union[str, List[str]]:
        """
        Extract content from a file.
        Must be implemented by subclasses.
        """
        pass

class MarkdownExtractor(BaseExtractor):
    """Extracts text from Markdown (.md) files."""

    @staticmethod
    def extract(file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logging.info(f"Successfully extracted Markdown file: {file_path}")
            return content
        except Exception as e:
            logging.error(f"Error reading Markdown file {file_path}: {e}")
            return ""

class PPTXExtractor(BaseExtractor):
    """Extracts slide-wise text from PowerPoint (.pptx) files."""

    @staticmethod
    def extract(file_path: str) -> List[str]:
        slides_text = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_contents = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_contents.append(shape.text)
                slide_text = "\n".join(slide_contents)
                slides_text.append(slide_text)
            logging.info(f"Successfully extracted PPTX file: {file_path}")
            return slides_text
        except Exception as e:
            logging.error(f"Error reading PPTX file {file_path}: {e}")
            return []

class FileProcessor:
    """
    Processes files within a root folder and extracts text content
    using appropriate extractor classes based on file extension.
    """

    def __init__(self, root_folder: str):
        self.root_folder = root_folder
        self.extractors: Dict[str, BaseExtractor] = {
            ".md": MarkdownExtractor,
            ".pptx": PPTXExtractor,
        }

    def process_files(self) -> Dict[str, Union[str, List[str]]]:
        """
        Recursively walk the root folder, extract content from supported files,
        and return a dictionary mapping file paths to their extracted content.
        """
        extracted_text: Dict[str, Union[str, List[str]]] = {}
        for dirpath, _, filenames in os.walk(self.root_folder):
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                current_file = os.path.join(dirpath, filename)
                extractor_class = self.extractors.get(ext)
                if extractor_class:
                    logging.info(f"Processing file: {current_file}")
                    extracted_text[current_file] = extractor_class.extract(current_file)
                else:
                    logging.debug(f"Skipping unsupported file type: {current_file}")
        return extracted_text
