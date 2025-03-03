import os
from pptx import Presentation

def extract_text_from_md(file_path):
    """Reads and returns the content of a markdown (.md) file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extracts and returns text from all slides in a PowerPoint (.pptx) file."""
    text_runs = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def process_files(root_folder):
    """
    Recursively walks through the root_folder, finds .md and .pptx files,
    and extracts their text content.
    Returns a dictionary where keys are file paths and values are the extracted text.
    Also prints the count of processed .md and .pptx files.
    """
    extracted_text = {}
    md_count = 0
    pptx_count = 0

    for dirpath, dirnames, filenames in os.walk(root_folder):
        for filename in filenames:
            file_path = os.path.join(dirpath, filename)
            if filename.lower().endswith(".md"):
                print(f"Processing Markdown file: {file_path}")
                extracted_text[file_path] = extract_text_from_md(file_path)
                md_count += 1
            elif filename.lower().endswith(".pptx"):
                print(f"Processing PowerPoint file: {file_path}")
                extracted_text[file_path] = extract_text_from_pptx(file_path)
                pptx_count += 1

    print("\nProcessing Summary:")
    print(f"Total Markdown files processed: {md_count}")
    print(f"Total PowerPoint files processed: {pptx_count}")

    return extracted_text

if __name__ == '__main__':
    result = process_files('./unprocessed_docs')

    # Output the results as a list of file names with their extracted text
    print("\nExtracted Text from Files:")
    for file_path, text in result.items():
        print("=" * 40)
        print(f"File: {file_path}")
        print("Extracted Text:")
        print(text)
        print("=" * 40)
