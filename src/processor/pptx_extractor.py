import logging
from pptx import Presentation
from typing import List
from .base_extractor import BaseExtractor

class PPTXExtractor(BaseExtractor):
    """Extracts slide-wise text from PowerPoint (.pptx) files."""

    @staticmethod
    def extract(file_path: str) -> List[str]:
        slides_text = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_contents = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_contents.append(shape.text)
                slide_text = "\n".join(slide_contents)
                slides_text.append(slide_text)
            logging.info(f"Successfully extracted PPTX file: {file_path}")
            return slides_text
        except Exception as e:
            logging.error(f"Error reading PPTX file {file_path}: {e}")
            return []
